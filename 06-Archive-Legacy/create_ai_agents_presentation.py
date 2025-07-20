import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Paths to images (relative to script location)
IMG_DIR = os.path.join('autowais', 'public', 'images', 'blog')

slide_images = [
    os.path.join(IMG_DIR, 'ai-automation.png'),         # Slide 1: AI concept
    os.path.join(IMG_DIR, 'workflow-automation.png'),  # Slide 2: Automation cycle
    os.path.join(IMG_DIR, 'remote-work-tools.png'),    # Slide 3: Future/benefits
]

titles = [
    'What Are AI Agents? 🤖',
    'How They Work & Examples You Know 🔄',
    "Why This Matters & What's Next 🚀",
]

contents = [
    [
        ('Key Differences from Regular AI:', True, False),
        ('Regular AI (ChatGPT):', True, False),
        ('You ask → It answers', False, False),
        ('Passive conversation', False, False),
        ('One response at a time', False, False),
        ('AI Agents:', True, False),
        ('You give a goal → It figures out how', False, False),
        ('Active problem-solving', False, False),
        ('Multiple actions until goal achieved', False, False),
        ('', False, False),
        ('Simple Analogy:', True, False),
        ('Regular AI = Smart encyclopedia that answers questions', False, False),
        ('AI Agent = Capable employee who handles entire projects', False, False),
        ('', False, False),
        ('The Magic Word: AUTONOMOUS', True, False),
        ('They can think, decide, act, and learn - all without constant human input!', False, False),
    ],
    [
        ('The 4-Step Cycle:', True, False),
        ('👁️ PERCEIVE - Gather information from environment', False, False),
        ('🧠 THINK - Analyze and plan actions', False, False),
        ('🤖 ACT - Take real-world actions', False, False),
        ('📚 LEARN - Improve from results', False, False),
        ('', False, False),
        ('This cycle repeats until the goal is achieved', False, False),
        ('', False, False),
        ('You Already Use AI Agents:', True, False),
        ('🗺️ Google Maps/Waze: Monitors traffic, reroutes, learns patterns', False, False),
        ('🏠 Smart Home Systems: Learns routines, adjusts automatically', False, False),
        ('🎵 Netflix/Spotify: Analyzes preferences, suggests content, improves over time', False, False),
    ],
    [
        ('Benefits for You:', True, False),
        ('✅ Save Time - Handle routine tasks automatically', False, False),
        ('✅ Never Forget - Continuous monitoring and reminders', False, False),
        ('✅ Work 24/7 - No breaks needed', False, False),
        ('✅ Personalized - Adapt to your preferences', False, False),
        ('✅ Cost Effective - Cheaper than human help', False, False),
        ('', False, False),
        ('Coming Soon:', True, False),
        ('Personal AI Assistants for your digital life', False, False),
        ('Business Automation for complex workflows', False, False),
        ('Smart Cities optimizing services', False, False),
        ('Healthcare Companions monitoring health', False, False),
        ('', False, False),
        ('Key Takeaway:', True, True),
        ("AI Agents aren't replacing human intelligence - they're amplifying it.", False, True),
        ('They handle routine tasks so you can focus on what matters most.', False, True),
        ('', False, False),
        ('The Bottom Line:', True, True),
        ("We're moving from tools that help us work → tools that work for us", False, True),
    ]
]

def add_image(slide, image_path):
    if os.path.exists(image_path) and image_path.lower().endswith(('.png', '.jpg', '.jpeg')):
        # Vertically center the image in the right column
        slide.shapes.add_picture(image_path, Inches(8.2), Inches(1.5), width=Inches(4.5), height=Inches(4.5))

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

for i in range(3):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    # Title
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.5), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = titles[i]
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    # Content (left column)
    content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7.2), Inches(5.5))
    content_frame = content_shape.text_frame
    for line, is_header, is_highlight in contents[i]:
        p = content_frame.add_paragraph()
        p.text = line
        if is_header and is_highlight:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black for highlight
            p.space_after = Pt(12)
        elif is_highlight:
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black for highlight
            p.space_after = Pt(10)
        elif is_header:
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.space_after = Pt(10)
        else:
            p.font.size = Pt(22)
            p.font.bold = False
            p.space_after = Pt(6)
    # Remove first empty paragraph
    if content_frame.paragraphs[0].text == '':
        content_frame._element.remove(content_frame.paragraphs[0]._element)
    # Image (right column)
    add_image(slide, slide_images[i])

# Save presentation
prs.save('AI_Agents_Overview.pptx')
print('Presentation created: AI_Agents_Overview.pptx') 